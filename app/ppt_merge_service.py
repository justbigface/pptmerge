from flask import Flask, request, send_file, jsonify
import io
import os
import tempfile
import requests
import copy
import contextlib
import urllib.parse
import time
import concurrent.futures
from pptx import Presentation

from pptx.opc.constants import RELATIONSHIP_TYPE as RT

app = Flask(__name__)


def clone_slide(src_slide, dest_prs):
    """Clone *src_slide* into *dest_prs* preserving shapes and media.

    This function takes a *source* slide and makes a deep‑copy of its XML
    into a *new* slide in *dest_prs*.  It also copies the relationships
    (images, charts, media, hyperlinks) so nothing goes missing.  It is
    not 100 % of PowerPoint features, but足以覆盖常见场景（文本框、图片、图表、形状、超链接）。
    """
    # 1) 选一个空白版式。大多数主题都有 index 6 = Blank，否则退回 0。
    blank_layout = dest_prs.slide_layouts[6] if len(dest_prs.slide_layouts) > 6 else dest_prs.slide_layouts[0]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # 2) 深拷贝 <p:spTree> 下的每个 shape 元素
    for shape in src_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        # _spTree.append 比 insert_element_before 简单且安全
        new_slide.shapes._spTree.append(new_el)

    # 3) 复制关系（图片 / 图表 / 媒体 / 超链接 / SmartArt / OLE 等）
    # 优先保留原 rId，若冲突则让 python-pptx 自动分配
    for rId, rel in src_slide.part.rels.items():
        try:
            # 尝试使用原 rId 复制关系
            new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.is_external, rId=rId)
        except KeyError:
            # 如果 rId 冲突，让 python-pptx 自动分配新的 rId
            new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.is_external)



def merge_presentations(streams):
    """Concatenate every slide of each PPTX *stream* into one *Presentation*."""
    merged = Presentation()

    # — 删除模板自动生成的空白首页（如果存在且为空白页）
    # 判断是否只有一页且该页的 spTree 下元素数≤1（更健壮，避免主题占位符）
    if len(merged.slides) == 1 and len(merged.slides[0].shapes._spTree) <= 1:
        merged.slides._sldIdLst.remove(merged.slides._sldIdLst[0])

    for s in streams:
        prs = Presentation(s)
        for slide in prs.slides:
            clone_slide(slide, merged)
    return merged


@app.route('/merge_pptx', methods=['POST'])
def merge_pptx():
    start_time = time.time()
    data = request.get_json(silent=True) or {}
    urls = data.get('urls', [])
    app.logger.info(f"Received merge request for URLs: {urls}")

    if not urls:
        app.logger.warning("No URLs provided in merge request.")
        return jsonify({'error': 'No urls provided'}), 400
    if len(urls) < 2:
        app.logger.warning("Less than two URLs provided in merge request.")
        return jsonify({'error': 'Need at least two URLs'}), 400

    temp_paths = []
    # 支持通过环境变量配置域名白名单和端口白名单，默认允许 justbigface.fun
    allowed_domains = [d.strip() for d in os.getenv('ALLOWED_HOSTS', '').split(',') if d.strip()]
    if not allowed_domains:
        allowed_domains = ['justbigface.fun']
    allowed_ports_env = os.getenv('ALLOWED_PORTS', '')
    if allowed_ports_env:
        allowed_ports = [int(p.strip()) for p in allowed_ports_env.split(',') if p.strip().isdigit()]
    else:
        allowed_ports = [80, 443, 9000, 9001]
    max_size_mb = int(os.getenv('MAX_SIZE_MB', '30'))
    max_size_bytes = max_size_mb * 1024 * 1024

    try:
        # 使用 ExitStack 管理文件流
        with contextlib.ExitStack() as stack:
            streams = []
            # 使用 requests.Session() 实现 HTTP keep-alive
            with requests.Session() as session:
                # —— 下载全部 PPTX ——
                def download_file(url, session, allowed_domains, allowed_ports, max_size_bytes):
                    parsed_url = urllib.parse.urlparse(url)
                    # SSRF防护：仅允许http/https，白名单域名，端口限制
                    if parsed_url.scheme not in ['http', 'https']:
                        safe_url = urllib.parse.urlsplit(url)._replace(query="[REDACTED]").geturl()
                        app.logger.warning(f"Invalid URL scheme: {safe_url}")
                        return {'error': f'Invalid URL scheme'}
                    # 域名白名单校验
                    if allowed_domains and parsed_url.hostname not in allowed_domains:
                        safe_url = urllib.parse.urlsplit(url)._replace(query="[REDACTED]").geturl()
                        app.logger.warning(f"Domain not in whitelist: {parsed_url.hostname} for URL {safe_url}")
                        return {'error': f'Domain not allowed: {parsed_url.hostname}'}
                    # 端口白名单校验
                    port = parsed_url.port or (443 if parsed_url.scheme == 'https' else 80)
                    if port not in allowed_ports:
                        safe_url = urllib.parse.urlsplit(url)._replace(query="[REDACTED]").geturl()
                        app.logger.warning(f"Port not allowed: {port} for URL {safe_url}")
                        return {'error': f'Port not allowed: {port}'}
                    try:
                        resp = session.get(url, stream=True, timeout=(5, 30))
                        resp.raise_for_status()
                        content_type = resp.headers.get('Content-Type', '')
                        if 'application/vnd.openxmlformats-officedocument.presentationml.presentation' not in content_type and not url.lower().endswith('.pptx'):
                            return {'error': f'Invalid Content-Type or file extension'}
                        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
                        downloaded_size = 0
                        for chunk in resp.iter_content(8192):
                            downloaded_size += len(chunk)
                            if downloaded_size > max_size_bytes:
                                tmp.close()
                                os.remove(tmp.name)
                                return {'error': f'File size exceeds limit ({max_size_bytes / 1024 / 1024} MB)'}
                            tmp.write(chunk)
                        tmp.close()
                        return {'success': tmp.name}
                    except Exception as e:
                        safe_url = urllib.parse.urlsplit(url)._replace(query="[REDACTED]").geturl()
                        app.logger.exception(f"Error downloading {safe_url}")
                        return {'error': f'Error downloading file'}

                # 并发度动态调整
                max_workers = min(5, len(urls))
                with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                    future_to_url = {executor.submit(download_file, url, session, allowed_domains, allowed_ports, max_size_bytes): url for url in urls}
                    for future in concurrent.futures.as_completed(future_to_url):
                        url = future_to_url[future]
                        try:
                            result = future.result()
                        except Exception as exc:
                            safe_url = urllib.parse.urlsplit(url)._replace(query="[REDACTED]").geturl()
                            app.logger.exception(f"Exception in downloading {safe_url}")
                            for p in temp_paths:
                                try:
                                    os.remove(p)
                                except OSError:
                                    pass
                            return jsonify({'error': f'Failed to download file'}), 400
                        if 'success' in result:
                            temp_paths.append(result['success'])
                        else:
                            for p in temp_paths:
                                try:
                                    os.remove(p)
                                except OSError:
                                    pass
                            return jsonify({'error': f'Failed to download file'}), 400

            # —— 打开文件流并合并 ——
            for p in temp_paths:
                streams.append(stack.enter_context(open(p, 'rb')))

            merged_prs = merge_presentations(streams)

            out = io.BytesIO()
            merged_prs.save(out)
            out.seek(0)
            end_time = time.time()
            duration = end_time - start_time
            app.logger.info(f"Successfully merged {len(urls)} files into {len(merged_prs.slides)} slides in {duration:.2f} seconds.")

            return send_file(
                out,
                as_attachment=True,
                download_name='merged.pptx',
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
    except Exception as exc:
        app.logger.error(f"Error merging PPTX files from {urls}: {exc}", exc_info=True)
        return jsonify({'error': 'An internal error occurred during merging.'}), 500
    finally:
        # 删除临时文件
        for p in temp_paths:
            try:
                os.remove(p)
            except OSError:
                app.logger.warning(f"Could not remove temporary file: {p}", exc_info=True)

def index():
    return 'PPT Merge Service is running!'


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080)
